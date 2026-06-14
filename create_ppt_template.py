"""
根据 docs/PPT_分工方案.md 生成 PPT 模板。
第1页：封面，第2-11页：空白占位，第12页：致谢/Q&A。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色方案 ──────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)   # 深蓝（标题）
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)  # 强调蓝
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF5)   # 浅灰背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x56, 0x65, 0x73)
MID_GRAY = RGBColor(0xAA, 0xAA, 0xAA)

SLIDE_W = Inches(13.333)  # 16:9 宽屏
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── 工具函数 ──────────────────────────────────────────
def add_bg(slide, color):
    """填充幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_BLUE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=DARK_BLUE, bold_first=False, line_spacing=1.5,
                  alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
        if bold_first and i == 0:
            p.font.bold = True
    return tf


def add_decorated_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    """添加装饰横线"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    """添加装饰横线（别名）"""
    return add_decorated_line(slide, left, top, width, color, height)


def add_slide_number(slide, num):
    """在右下角添加页码"""
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
                str(num), font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════
# 第 1 页：封面
# ═══════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
add_bg(slide1, DARK_BLUE)

# 顶部装饰条
add_decorated_line(slide1, Inches(0), Inches(0), SLIDE_W, ACCENT_BLUE, Pt(5))

# 中文标题
add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.8),
            "基于强化学习的\n人员配置-生产调度协同优化",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 英文副标题
add_textbox(slide1, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
            "Reinforcement Learning for Personnel Allocation & Production Scheduling",
            font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# 分隔线
add_decorated_line(slide1, Inches(5), Inches(4.5), Inches(3.3), ACCENT_BLUE, Pt(2))

# 组员 / 课程 / 日期
add_multiline(slide1, Inches(3.5), Inches(5.0), Inches(6.3), Inches(1.5),
              ["组员：__________  /  __________  /  __________  /  __________  /  __________",
               "课程：________________________________________",
               "日期：2026.06"],
              font_size=16, color=RGBColor(0xCC, 0xD1, 0xD9), alignment=PP_ALIGN.CENTER,
              line_spacing=2.0)

add_slide_number(slide1, 1)


# ═══════════════════════════════════════════════════════
# 第 2-11 页：空白占位页
# ═══════════════════════════════════════════════════════
placeholder_pages = [
    (2,  "问题背景",          "Person A — 问题引入"),
    (3,  "方法总览",          "Person A — 方法总览"),
    (4,  "SMDP 环境建模",     "Person B — 方法细节"),
    (5,  "DQN & PPO 算法设计","Person B — 方法细节"),
    (6,  "核心创新：ScaleInv 架构", "Person B — 方法细节"),
    (7,  "实验设计",          "Person C — 实验设计"),
    (8,  "消融实验结果",      "Person C — 实验设计"),
    (9,  "最终测试结果",      "Person D — 结果展示"),
    (10, "可视化分析",        "Person D — 结果展示"),
    (11, "结论与核心贡献",    "Person E — 总结展望"),
]

for page_num, title, subtitle in placeholder_pages:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # 顶部色条
    add_decorated_line(slide, Inches(0), Inches(0), SLIDE_W, DARK_BLUE, Pt(4))

    # 页标题
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
                title, font_size=32, color=DARK_BLUE, bold=True)

    # 副标题（灰色，说明谁负责）
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.4),
                subtitle, font_size=13, color=MID_GRAY)

    # 分隔线
    add_decorated_line(slide, Inches(0.8), Inches(1.6), Inches(11.7), ACCENT_BLUE, Pt(2))

    # 中央占位提示（浅色虚线框效果）
    placeholder_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(2.5), Inches(10.3), Inches(4.2)
    )
    placeholder_shape.fill.solid()
    placeholder_shape.fill.fore_color.rgb = LIGHT_GRAY
    placeholder_shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    placeholder_shape.line.width = Pt(1)

    tf = placeholder_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"第 {page_num} 页\n{title}"
    p.font.size = Pt(22)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "（待填写内容）"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(100)

    add_slide_number(slide, page_num)


# ═══════════════════════════════════════════════════════
# 第 12 页：致谢 / Q&A
# ═══════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12, DARK_BLUE)

# 顶部装饰条
add_decorated_line(slide12, Inches(0), Inches(0), SLIDE_W, ACCENT_BLUE, Pt(5))

# 感谢
add_textbox(slide12, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
            "感谢聆听！", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

# 欢迎提问
add_textbox(slide12, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.6),
            "欢迎提问 & 讨论", font_size=22, color=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER)

# 分隔线
add_decorated_line(slide12, Inches(5), Inches(3.6), Inches(3.3), ACCENT_BLUE, Pt(2))

# 资源信息
add_multiline(slide12, Inches(3.0), Inches(4.1), Inches(7.3), Inches(1.5),
              ["代码仓库：________________________________________",
               "报告文档：________________________________________",
               "模型 & 可视化：____________________________________"],
              font_size=14, color=RGBColor(0xCC, 0xD1, 0xD9),
              alignment=PP_ALIGN.CENTER, line_spacing=2.0)

# 分工信息
add_multiline(slide12, Inches(2.0), Inches(5.7), Inches(9.3), Inches(1.6),
              ["分工：",
               "Person A：问题引入 & 方法总览    |    Person B：环境建模 & 算法设计 & 核心创新",
               "Person C：实验设计 & 消融验证      |    Person D：最终结果 & 可视化分析",
               "Person E：结论总结 & 致谢"],
              font_size=11, color=RGBColor(0xAA, 0xB2, 0xBF),
              alignment=PP_ALIGN.CENTER, line_spacing=1.8)

add_slide_number(slide12, 12)

# ── 保存 ──────────────────────────────────────────────
output_path = "docs/PPT_模板.pptx"
prs.save(output_path)
print(f"PPT template saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
print(f"   Slide 1: Title page")
print(f"   Slides 2-11: Placeholder pages")
print(f"   Slide 12: Thank you / Q&A")
